#!/usr/bin/env python3
"""Build `docs/deck.pptx` from the same content as `docs/deck.html`.

Native slides with real text boxes, not a PDF rendered to images, so the deck
can be edited, re-ordered and presented in PowerPoint, Keynote or Google
Slides. The only bitmap is the architecture diagram, because PowerPoint's SVG
support is not dependable enough to rely on for the one slide that carries the
argument.

`deck.html` remains what gets presented and what prints to `deck.pdf`. This
file mirrors it. If the two disagree, this is the bug.

    uv run --with python-pptx python scripts/build_deck_pptx.py
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

DOCS = Path(__file__).resolve().parent.parent / "docs"

INK = RGBColor(0x0F, 0x17, 0x2A)
INK2 = RGBColor(0x33, 0x41, 0x55)
INK3 = RGBColor(0x64, 0x74, 0x8B)
PAGE = RGBColor(0xFB, 0xFA, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LINE = RGBColor(0xE2, 0xE8, 0xF0)
MODEL = RGBColor(0x4F, 0x46, 0xE5)
CODE = RGBColor(0x0F, 0x76, 0x6E)
WARN = RGBColor(0xB4, 0x53, 0x09)

FONT = "Segoe UI"
MONO = "Consolas"

W = Inches(13.333)
H = Inches(7.5)
MARGIN = Inches(0.72)


def deck() -> Presentation:
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    return prs


def slide(prs: Presentation):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = PAGE
    return s


def text(
    s,
    x,
    y,
    w,
    h,
    runs,
    *,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    spacing=1.0,
):
    """`runs` is a list of (text, size, bold, colour) or a list of such lists.

    A list of lists makes one paragraph each, which is how the bullets and the
    stacked headings are built.
    """
    box = s.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

    paragraphs = runs if runs and isinstance(runs[0], list) else [runs]
    for index, para_runs in enumerate(paragraphs):
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        if index:
            p.space_before = Pt(10)
        for content, size, bold, colour in para_runs:
            r = p.add_run()
            r.text = content
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = colour
            r.font.name = FONT
    return box


def card(s, x, y, w, h, *, fill=WHITE, border=LINE):
    shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.adjustments[0] = 0.06
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = border
    shape.line.width = Pt(1)
    shape.shadow.inherit = False
    shape.text_frame.text = ""
    return shape


def rule(s, x, y, *, width=Inches(0.95)):
    bar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, width, Pt(4))
    bar.fill.solid()
    bar.fill.fore_color.rgb = MODEL
    bar.line.fill.background()
    bar.shadow.inherit = False
    return bar


def heading(s, eyebrow, title, *, title_size=34):
    text(s, MARGIN, Inches(0.62), Inches(9), Inches(0.3),
         [(eyebrow.upper(), 11, True, MODEL)])
    rule(s, MARGIN, Inches(0.98))
    text(s, MARGIN, Inches(1.22), Inches(11.9), Inches(1.2),
         [(title, title_size, True, INK)], spacing=1.05)


def foot(s, note):
    text(s, MARGIN, Inches(6.66), Inches(11.4), Inches(0.4),
         [(note, 11.5, False, INK3)])


def page_number(s, n):
    text(s, Inches(12.2), Inches(6.86), Inches(0.6), Inches(0.3),
         [(str(n), 10.5, False, INK3)], align=PP_ALIGN.RIGHT)


def bullets(s, x, y, w, h, items, *, size=13.5, colour=INK2):
    paras = []
    for item in items:
        runs = [("•   ", size, False, colour)]
        runs.extend(
            (chunk, size, bold, INK if bold else colour)
            for chunk, bold in item
        )
        paras.append(runs)
    return text(s, x, y, w, h, paras, spacing=1.25)


def build() -> Presentation:
    prs = deck()

    # ---------------------------------------------------------------- 1
    s = slide(prs)
    s.shapes.add_picture(
        str(DOCS / "media" / "logo.png"), Inches(4.9), Inches(1.5), width=Inches(3.5)
    )
    text(s, MARGIN, Inches(3.0), Inches(11.9), Inches(1.5),
         [[("A crew desk advisor", 40, True, INK)],
          [("that never guesses", 40, True, INK)]],
         align=PP_ALIGN.CENTER, spacing=1.05)
    text(s, Inches(2.4), Inches(4.72), Inches(8.5), Inches(1.0),
         [("The model plans and explains. Deterministic code computes. "
           "A guard checks every figure against what the tools returned.",
           15, False, INK2)],
         align=PP_ALIGN.CENTER, spacing=1.35)
    text(s, MARGIN, Inches(5.62), Inches(11.9), Inches(0.4),
         [("extroc-jpkcqxtlma-uc.a.run.app", 14, True, MODEL)],
         align=PP_ALIGN.CENTER)
    text(s, MARGIN, Inches(6.86), Inches(11.9), Inches(0.3),
         [("dCortex, Agentic Crew Ops Advisor", 11, False, INK3)],
         align=PP_ALIGN.CENTER)

    # ---------------------------------------------------------------- 2
    s = slide(prs)
    heading(s, "The problem", "The bottleneck is not detecting that something broke")
    text(s, MARGIN, Inches(2.16), Inches(10.4), Inches(0.8),
         [("It is reasoning correctly, and fast, about what follows. A sick call at "
           "05:00 breaks a pairing, strands downstream legs, and every fix creates "
           "the next problem.", 14, False, INK2)], spacing=1.35)
    cards = [
        ("Fragmented",
         "One answer spans rosters, duty clocks, schedules, reserves, "
         "qualifications and the rulebook."),
        ("Consequence blind",
         "The broken flight is obvious. The four that break next are not."),
        ("Exact, not approximate",
         "Legality is arithmetic against a rulebook. An approximate answer "
         "is a violation."),
    ]
    for i, (title, body) in enumerate(cards):
        x = MARGIN + Emu(int(i * Inches(4.02)))
        card(s, x, Inches(3.3), Inches(3.75), Inches(1.95))
        text(s, x + Inches(0.28), Inches(3.55), Inches(3.2), Inches(0.35),
             [(title, 15, True, INK)])
        text(s, x + Inches(0.28), Inches(4.0), Inches(3.2), Inches(1.1),
             [(body, 12, False, INK2)], spacing=1.3)
    foot(s, "Today that reasoning lives in one experienced controller's head. "
            "It degrades exactly when it matters most.")
    page_number(s, 2)

    # ---------------------------------------------------------------- 3
    s = slide(prs)
    text(s, MARGIN, Inches(0.62), Inches(9), Inches(0.3),
         [("THE QUESTION ACTUALLY BEING ASKED", 11, True, MODEL)])
    rule(s, MARGIN, Inches(0.98))
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN, Inches(2.0), Pt(4), Inches(1.55))
    bar.fill.solid()
    bar.fill.fore_color.rgb = MODEL
    bar.line.fill.background()
    bar.shadow.inherit = False
    text(s, MARGIN + Inches(0.3), Inches(2.0), Inches(9.6), Inches(1.6),
         [("What should the language model do, what should deterministic code do, "
           "and how do you compose them into a system that is both conversational "
           "and correct?", 21, True, INK)], spacing=1.3)
    text(s, MARGIN, Inches(4.25), Inches(9.6), Inches(1.2),
         [("Put the data in the prompt and let the model answer, and Tier 1 works. "
           "Tier 2 and Tier 3 fail. A model that approximates a duty hour "
           "calculation produces answers that are fluent, confident and wrong.",
           14, False, INK2)], spacing=1.35)
    text(s, MARGIN, Inches(5.45), Inches(9.6), Inches(0.4),
         [("Operationally, that is worse than no answer.", 15, True, INK)])
    page_number(s, 3)

    # ---------------------------------------------------------------- 4
    s = slide(prs)
    text(s, MARGIN, Inches(0.62), Inches(9), Inches(0.3),
         [("OUR ANSWER", 11, True, MODEL)])
    rule(s, MARGIN, Inches(0.98))
    text(s, MARGIN, Inches(1.22), Inches(11.9), Inches(1.1),
         [[("The model plans and explains.", 32, True, INK)],
          [("It never produces a fact.", 32, True, INK)]], spacing=1.06)
    panels = [
        (MARGIN, "THE MODEL MAY", MODEL, RGBColor(0xEE, 0xF2, 0xFF), [
            [("Decide which tools to call, with what arguments, in what order", False)],
            [("Decide when it has enough, and when it must decline", False)],
            [("Phrase the result for someone with a radio in one hand", False)],
        ]),
        (Inches(7.0), "ONLY CODE MAY", CODE, RGBColor(0xEC, 0xFD, 0xF5), [
            [("Compute any number, and show the arithmetic that produced it", False)],
            [("Decide that an assignment is legal or illegal", False)],
            [("Verify the drafted answer before a controller sees it", False)],
        ]),
    ]
    for x, label, colour, tint, items in panels:
        card(s, x, Inches(3.05), Inches(5.6), Inches(2.75))
        chip = s.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.3), Inches(3.3),
            Inches(1.85), Inches(0.32))
        chip.adjustments[0] = 0.5
        chip.fill.solid()
        chip.fill.fore_color.rgb = tint
        chip.line.fill.background()
        chip.shadow.inherit = False
        text(s, x + Inches(0.3), Inches(3.36), Inches(1.85), Inches(0.25),
             [(label, 9.5, True, colour)], align=PP_ALIGN.CENTER)
        bullets(s, x + Inches(0.3), Inches(3.85), Inches(5.0), Inches(1.7), items,
                size=12.5)
    foot(s, "A turn crosses that boundary four times. The next slide is the shape of it.")
    page_number(s, 4)

    # ---------------------------------------------------------------- 5
    s = slide(prs)
    text(s, MARGIN, Inches(0.5), Inches(9), Inches(0.3),
         [("ARCHITECTURE", 11, True, MODEL)])
    s.shapes.add_picture(
        str(DOCS / "media" / "architecture.png"),
        Inches(0.42), Inches(1.05), width=Inches(12.5))
    page_number(s, 5)

    # ---------------------------------------------------------------- 6
    s = slide(prs)
    heading(s, "Why it holds", "Enforced, not requested")
    text(s, MARGIN, Inches(2.16), Inches(10.4), Inches(0.5),
         [("A guarantee written into a prompt is a request. These are the three "
           "places the boundary is actually held.", 14, False, INK2)], spacing=1.35)
    rows = [
        ("A build test walks the import graph of domain, rules, ops, store, tools "
         "and verify, and fails if a model client is reachable",
         "Arithmetic drifting into the model's half"),
        ("The verifier rejects any atom in the prose that no tool emitted as a "
         "Fact this turn",
         "The model stating a plausible number nobody computed"),
        ("Graph edges require a legality result before any verdict, and a cover "
         "search before any recommendation",
         "The model inferring a verdict from context"),
    ]
    table = s.shapes.add_table(
        4, 2, MARGIN, Inches(2.95), Inches(11.9), Inches(2.6)).table
    table.columns[0].width = Inches(7.3)
    table.columns[1].width = Inches(4.6)
    # Rows size themselves to content in PowerPoint, but the header keeps the
    # generous default unless it is told otherwise, which left a band of empty
    # cell above the first mechanism.
    table.rows[0].height = Inches(0.34)
    for col, label in enumerate(("MECHANISM", "WHAT IT STOPS")):
        cell = table.cell(0, col)
        cell.text = label
        p = cell.text_frame.paragraphs[0]
        p.runs[0].font.size = Pt(10.5)
        p.runs[0].font.bold = True
        p.runs[0].font.color.rgb = INK3
        p.runs[0].font.name = FONT
        cell.fill.solid()
        cell.fill.fore_color.rgb = PAGE
    for r, (mech, stops) in enumerate(rows, start=1):
        for col, body in enumerate((mech, stops)):
            cell = table.cell(r, col)
            cell.text = body
            cell.fill.solid()
            cell.fill.fore_color.rgb = PAGE
            p = cell.text_frame.paragraphs[0]
            p.runs[0].font.size = Pt(12)
            p.runs[0].font.color.rgb = INK if col else INK2
            p.runs[0].font.name = FONT
    foot(s, "The third matters most. A guarantee written as a graph edge is a guarantee.")
    page_number(s, 6)

    # ---------------------------------------------------------------- 7
    s = slide(prs)
    heading(s, "Explainability", "Every figure carries its own working")
    block = card(s, MARGIN, Inches(2.35), Inches(5.9), Inches(3.15),
                 fill=INK, border=INK)
    block.text_frame.word_wrap = True
    code_lines = [
        "Fact(",
        '  key="C-2087.duty_7d.projected",',
        "  value=61.33,",
        '  unit="hours",',
        "  provenance=COMPUTED,",
        '  source="crewops.rules.duty.window",',
        '  derivation="51.83h prior + 9.50h from',
        '    P-2291 = 61.33h against a 60.00h',
        '    limit, over by 1.33h",',
        ")",
    ]
    tf = block.text_frame
    tf.margin_left = Inches(0.3)
    tf.margin_top = Inches(0.25)
    for i, line in enumerate(code_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.25
        r = p.add_run()
        r.text = line
        r.font.size = Pt(11.5)
        r.font.name = MONO
        r.font.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
    text(s, Inches(7.1), Inches(2.4), Inches(5.5), Inches(3.2),
         [[("derivation is the point. ", 14, True, INK),
           ("A controller about to move a crew member and sign their name to it "
            "does not want to be told the answer. They want to check it and argue "
            "with it.", 14, False, INK2)],
          [("It is also what makes the verifier possible. If a number can appear "
            "in an answer, a tool must have emitted a Fact for it.", 14, False, INK2)],
          [("When verification fails, the fix is to add the missing fact to the "
            "tool. It is never to relax the check.", 14, True, INK)]],
         spacing=1.35)
    page_number(s, 7)

    # ---------------------------------------------------------------- 8
    s = slide(prs)
    heading(s, "Measured", "44 cases: 38 shipped questions, 6 worked scenarios")
    data = [
        ("Tier", "Cases", "Correct", "Abstained", "Wrong", "Accuracy when answered"),
        ("Tier 1, lookup", "16", "15", "1", "0", "100%"),
        ("Tier 2, consequence", "14", "10", "4", "0", "100%"),
        ("Tier 3, recommendation", "14", "10", "4", "0", "100%"),
        ("Overall", "44", "35", "9", "0", "100%"),
    ]
    table = s.shapes.add_table(5, 6, MARGIN, Inches(2.3), Inches(11.9), Inches(2.1)).table
    widths = [3.5, 1.3, 1.4, 1.6, 1.2, 2.9]
    for i, inches in enumerate(widths):
        table.columns[i].width = Inches(inches)
    for r, row in enumerate(data):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = value
            cell.fill.solid()
            cell.fill.fore_color.rgb = PAGE
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.RIGHT
            run = p.runs[0]
            run.font.name = FONT
            if r == 0:
                run.font.size = Pt(10.5)
                run.font.bold = True
                run.font.color.rgb = INK3
            else:
                run.font.size = Pt(13)
                run.font.bold = r == 4
                run.font.color.rgb = INK if r == 4 else INK2
    stats = [
        ("0", "wrong answers, and no verdict inversions", CODE),
        ("35", "of 44 answers passed grounding verification", INK),
        ("122ms", "p95 on the deterministic path", INK),
    ]
    for i, (big, label, colour) in enumerate(stats):
        x = MARGIN + Emu(int(i * Inches(4.02)))
        text(s, x, Inches(4.85), Inches(3.7), Inches(0.8),
             [(big, 42, True, colour)])
        text(s, x, Inches(5.72), Inches(3.6), Inches(0.6),
             [(label, 11.5, False, INK3)], spacing=1.25)
    foot(s, "Reproduce with make eval. Deterministic path, which runs with no API key.")
    page_number(s, 8)

    # ---------------------------------------------------------------- 9
    s = slide(prs)
    heading(s, "The design choice we would defend hardest",
            "Abstention is a feature, and it is scored as one")
    text(s, MARGIN, Inches(2.3), Inches(10.6), Inches(0.9),
         [("The evaluation harness counts abstentions separately from wrong answers "
           "and never treats one as a failure. A grader that scored refusal as "
           "failure would push the system toward confident guessing, which is the "
           "exact failure mode the brief warns about.", 14, False, INK2)], spacing=1.35)
    panels = [
        (MARGIN, "CORRECT OUTCOME", CODE, RGBColor(0xEC, 0xFD, 0xF5),
         "\"I cannot answer that reliably.\" Then: what was missing, what was "
         "established anyway, and three questions that would work instead."),
        (Inches(7.0), "NEVER ACCEPTABLE", WARN, RGBColor(0xFF, 0xFB, 0xEB),
         "A fluent, confident figure nobody computed. The verifier rejects it "
         "rather than trimming it, so it cannot reach the screen."),
    ]
    for x, label, colour, tint, body in panels:
        card(s, x, Inches(3.65), Inches(5.6), Inches(1.85))
        chip = s.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.3), Inches(3.9),
            Inches(2.0), Inches(0.32))
        chip.adjustments[0] = 0.5
        chip.fill.solid()
        chip.fill.fore_color.rgb = tint
        chip.line.fill.background()
        chip.shadow.inherit = False
        text(s, x + Inches(0.3), Inches(3.96), Inches(2.0), Inches(0.25),
             [(label, 9.5, True, colour)], align=PP_ALIGN.CENTER)
        text(s, x + Inches(0.3), Inches(4.45), Inches(5.0), Inches(0.9),
             [(body, 12.5, False, INK2)], spacing=1.3)
    foot(s, "9 of 44 cases abstained. Every one of them declined honestly rather "
            "than guessing.")
    page_number(s, 9)

    # ---------------------------------------------------------------- 10
    s = slide(prs)
    heading(s, "Beyond the brief", "Voice is a peripheral, not a second brain")
    text(s, MARGIN, Inches(2.3), Inches(10.6), Inches(0.8),
         [("Speech in becomes a transcript. The transcript goes to the same "
           "endpoint, through the same tools, the same seven rules and the same "
           "verifier. Speech out reads prose the verifier has already passed.",
           14, False, INK2)], spacing=1.35)
    bullets(s, MARGIN, Inches(3.5), Inches(11.4), Inches(2.2), [
        [("Nothing under agent/voice/ imports a model client. No speech provider "
          "ever sees the dataset.", False)],
        [("A draft that fails verification is ", False), ("never spoken", True),
         (": the selector returns nothing at all unless the reply is verified "
          "or repaired.", False)],
        [("Hands free, it reads the verified headline first and offers the detail, "
          "so a controller can keep their hands on the desk.", False)],
    ], size=14)
    foot(s, "Also shipped: proactive alerting, multi turn memory, drafted crew "
            "notifications, chained disruptions.")
    page_number(s, 10)

    # ---------------------------------------------------------------- 11
    s = slide(prs)
    heading(s, "Honest limits", "What breaks, and how badly")
    text(s, MARGIN, Inches(2.22), Inches(10.6), Inches(0.6),
         [("Every failure is graded safe or unsafe. A safe failure declines. An "
           "unsafe failure answers wrongly. We treat the second as a different "
           "class of problem.", 13.5, False, INK2)], spacing=1.35)
    bullets(s, MARGIN, Inches(3.15), Inches(11.4), Inches(3.2), [
        [("Agent mode is not reproducible.", True),
         (" Three identical passes scored 16, 15 and 16. We quote a range, "
          "never a single number.", False)],
        [("Nine cases abstain", True),
         (" that a stronger router would answer. Each is a routing gap, not a "
          "reasoning error.", False)],
        [("The eighth rule problem.", True),
         (" Answer keys exclude candidates for reasons the rulebook does not "
          "cover, most importantly double booking. Modelling it as a RULE- id "
          "would misrepresent the rulebook, so it is carried as an operational "
          "feasibility issue instead: blocking, but honestly labelled.", False)],
        [("Grounding is per atom, not per claim.", True),
         (" It catches an invented figure. It would not catch a correctly quoted "
          "figure used in a wrong sentence.", False)],
    ], size=13)
    foot(s, "Full analysis, failure by failure, in docs/FAILURE-ANALYSIS.md")
    page_number(s, 11)

    # ---------------------------------------------------------------- 12
    s = slide(prs)
    text(s, MARGIN, Inches(0.62), Inches(9), Inches(0.3),
         [("IF THIS WERE REAL", 11, True, MODEL)])
    rule(s, MARGIN, Inches(0.98))
    panels = [
        ("Impact",
         "The work is the cross referencing, not the decision. Collapsing that "
         "from minutes to seconds is the value, and the ranked options make the "
         "decision reviewable afterwards."),
        ("Scale",
         "WorldState is loaded once and immutable, with a SQLite projection for "
         "lookups. The rules engine is pure arithmetic over typed records, so it "
         "scales with crew count, not with prompt size."),
        ("Crew PII",
         "The model never needs identity. It plans over ids and receives Facts. "
         "Names can stay behind the tool boundary and be joined at render time, "
         "so PII need never enter a prompt."),
    ]
    for i, (title, body) in enumerate(panels):
        x = MARGIN + Emu(int(i * Inches(4.02)))
        card(s, x, Inches(2.1), Inches(3.75), Inches(3.3))
        text(s, x + Inches(0.3), Inches(2.4), Inches(3.2), Inches(0.4),
             [(title, 17, True, INK)])
        text(s, x + Inches(0.3), Inches(2.95), Inches(3.15), Inches(2.2),
             [(body, 12.5, False, INK2)], spacing=1.35)
    foot(s, "Reasoning and arithmetic behind each of these in docs/PRODUCTION.md")
    page_number(s, 12)

    # ---------------------------------------------------------------- 13
    s = slide(prs)
    text(s, MARGIN, Inches(2.2), Inches(11.9), Inches(0.3),
         [("LIVE", 11, True, MODEL)], align=PP_ALIGN.CENTER)
    text(s, MARGIN, Inches(2.65), Inches(11.9), Inches(0.7),
         [("Ask it something it cannot answer", 34, True, INK)],
         align=PP_ALIGN.CENTER)
    text(s, Inches(2.9), Inches(3.6), Inches(7.5), Inches(1.2),
         [("That is the demo we would run first. A system that says \"I cannot "
           "answer that reliably, and here is what was missing\" is the one worth "
           "having at 06:00 on a bad day.", 15, False, INK2)],
         align=PP_ALIGN.CENTER, spacing=1.4)
    text(s, MARGIN, Inches(4.95), Inches(11.9), Inches(0.4),
         [("extroc-jpkcqxtlma-uc.a.run.app", 16, True, MODEL)],
         align=PP_ALIGN.CENTER)
    text(s, MARGIN, Inches(5.45), Inches(11.9), Inches(0.4),
         [("Runs with no API key. The deterministic path answers through the same "
           "tools, rules and verifier.", 11.5, False, INK3)],
         align=PP_ALIGN.CENTER)
    page_number(s, 13)

    return prs


if __name__ == "__main__":
    out = DOCS / "deck.pptx"
    build().save(str(out))
    print(f"wrote {out}")
